"""Genera PPT profesional del análisis nacional y por regiones."""
from __future__ import annotations

from pathlib import Path

import matplotlib.pyplot as plt
import matplotlib.ticker as mticker
import numpy as np
import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

BASE = Path(__file__).parent
CSV = BASE / 'trayectorias_imposibles_por_region.csv'
OUT_PPT = BASE / 'presentacion_trayectorias.pptx'
FIG_DIR = BASE / '_ppt_figs'

# Paleta presentación
NAVY = '#0B2545'
TEAL = '#1B7A6E'
ORANGE = '#C45C26'
GRAY = '#5A6A7A'
LIGHT = '#F4F6F8'
COLOR_1A = '#1B6B9A'
COLOR_2A = '#C45C26'
COLOR_REV = '#6B7C8C'
TIPO_COLOR = {
    'Bosque aislado 1 año': COLOR_1A,
    'Bosque aislado 2 años': COLOR_2A,
    'Revisar patrón': COLOR_REV,
}


def _rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip('#')
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _fmt_k(x, _pos=None):
    if abs(x) >= 1_000_000:
        return f'{x / 1_000_000:.1f}M'
    if abs(x) >= 1_000:
        return f'{x / 1_000:.0f}k'
    return f'{int(x)}'


def _style():
    plt.rcParams.update({
        'figure.facecolor': 'white',
        'axes.facecolor': 'white',
        'axes.edgecolor': '#D0D7DE',
        'axes.grid': True,
        'grid.color': '#EEF1F4',
        'grid.linewidth': 0.8,
        'axes.axisbelow': True,
        'font.size': 11,
        'axes.titlesize': 14,
        'axes.titleweight': 'bold',
        'axes.labelcolor': NAVY,
        'xtick.color': GRAY,
        'ytick.color': GRAY,
        'text.color': NAVY,
    })


def _save(fig, name: str) -> Path:
    FIG_DIR.mkdir(exist_ok=True)
    path = FIG_DIR / f'{name}.png'
    fig.savefig(path, dpi=200, bbox_inches='tight', facecolor='white')
    plt.close(fig)
    return path


def cargar():
    df = pd.read_csv(CSV, encoding='utf-8-sig')
    df['region_id'] = df['region_id'].astype(str)
    if 'trayectoria' not in df.columns:
        if 'clases' in df.columns:
            df['trayectoria'] = df['clases']
        else:
            def cod(k):
                k = int(float(k))
                return f'{k // 1_000_000}-{(k % 1_000_000) // 10_000}-{(k % 10_000) // 100}-{k % 100}'
            df['trayectoria'] = df['codigo'].apply(cod)

    nac = (
        df.groupby(['codigo', 'trayectoria', 'tipo'], as_index=False)
        .agg(pixeles=('pixeles', 'sum'), n_regiones=('region_id', 'nunique'))
        .sort_values('pixeles', ascending=False)
        .reset_index(drop=True)
    )
    nac['pct_pixeles'] = 100 * nac['pixeles'] / nac['pixeles'].sum()

    por_region = (
        df.groupby('region_id', as_index=False)
        .agg(pixeles=('pixeles', 'sum'), patrones=('codigo', 'nunique'))
        .sort_values('pixeles', ascending=False)
    )
    por_tipo = (
        nac.groupby('tipo', as_index=False)['pixeles']
        .sum()
        .sort_values('pixeles', ascending=False)
    )
    return df, nac, por_region, por_tipo


def fig_tipo_donut(por_tipo, total_px):
    fig, ax = plt.subplots(figsize=(8, 5.2))
    colors = [TIPO_COLOR.get(t, GRAY) for t in por_tipo['tipo']]
    labels = [f"{t.split(' ')[-1] if 'año' in t else t}\n{100 * p / total_px:.0f}%"
              for t, p in zip(por_tipo['tipo'], por_tipo['pixeles'])]
    # etiquetas más claras
    labels = []
    for t, p in zip(por_tipo['tipo'], por_tipo['pixeles']):
        short = {'Bosque aislado 1 año': '1 año', 'Bosque aislado 2 años': '2 años'}.get(t, 'Revisar')
        labels.append(f'{short}\n{100 * p / total_px:.0f}%')
    wedges, _ = ax.pie(
        por_tipo['pixeles'].values, colors=colors, startangle=90,
        wedgeprops=dict(width=0.42, edgecolor='white', linewidth=2.5),
    )
    ax.legend(wedges, labels, loc='center', frameon=False, fontsize=11)
    ax.text(0, 0, f'{total_px / 1e6:.1f}M\npíxeles', ha='center', va='center',
            fontsize=16, fontweight='bold', color=NAVY)
    return _save(fig, '01_tipo')


def fig_top_tray(nac, total_px, n=15):
    top = nac.head(n)
    fig, ax = plt.subplots(figsize=(10, 6.5))
    y = np.arange(len(top))
    ax.barh(y, top['pixeles'], color=[TIPO_COLOR.get(t, GRAY) for t in top['tipo']],
            height=0.7, edgecolor='white')
    ax.set_yticks(y, top['trayectoria'], family='monospace', fontsize=10)
    ax.invert_yaxis()
    ax.xaxis.set_major_formatter(mticker.FuncFormatter(_fmt_k))
    ax.set_xlabel('Píxeles')
    for i, (_, r) in enumerate(top.iterrows()):
        ax.text(r['pixeles'] + total_px * 0.003, i, f"{r['pct_pixeles']:.1f}%",
                va='center', fontsize=8, color=GRAY)
    ax.set_xlim(0, top['pixeles'].max() * 1.16)
    fig.subplots_adjust(left=0.18)
    return _save(fig, '02_top_tray')


def fig_concentracion(nac, total_px):
    cum = nac['pixeles'].cumsum() / total_px
    n80 = int(np.searchsorted(cum.values, 0.8) + 1)
    fig, ax = plt.subplots(figsize=(9, 4.8))
    x = np.arange(1, len(cum) + 1)
    ax.fill_between(x, cum.values, color=COLOR_1A, alpha=0.12)
    ax.plot(x, cum.values, color=COLOR_1A, lw=2.5)
    ax.axhline(0.8, color=ORANGE, ls='--', lw=1.3)
    ax.axvline(n80, color=ORANGE, ls='--', lw=1.0, alpha=0.8)
    ax.set_ylim(0, 1.02)
    ax.set_xlabel('Patrones (mayor → menor)')
    ax.set_ylabel('Fracción acumulada')
    ax.annotate(f'{n80} patrones ≈ 80%', xy=(n80, 0.8),
                xytext=(n80 + max(40, len(nac) * 0.08), 0.55),
                fontsize=11, color=ORANGE,
                arrowprops=dict(arrowstyle='->', color=ORANGE))
    return _save(fig, '03_concentracion'), n80


def fig_extension(nac):
    fig, ax = plt.subplots(figsize=(9, 5.2))
    for tipo, g in nac.groupby('tipo'):
        ax.scatter(
            g['n_regiones'], g['pixeles'],
            s=np.clip(35 + g['pct_pixeles'] * 6, 18, 220),
            c=TIPO_COLOR.get(tipo, GRAY), alpha=0.75,
            edgecolors='white', linewidths=0.5,
            label={'Bosque aislado 1 año': '1 año', 'Bosque aislado 2 años': '2 años'}.get(tipo, 'Revisar'),
        )
    for _, r in nac.head(6).iterrows():
        ax.annotate(r['trayectoria'], (r['n_regiones'], r['pixeles']),
                    textcoords='offset points', xytext=(5, 3),
                    fontsize=8, family='monospace', color=NAVY)
    ax.set_yscale('log')
    ax.yaxis.set_major_formatter(mticker.FuncFormatter(_fmt_k))
    ax.set_xlabel('N regiones')
    ax.set_ylabel('Píxeles')
    ax.legend(frameon=False, fontsize=10)
    return _save(fig, '04_extension')


def fig_regiones_px_cambios(por_region, top=12):
    fig, axes = plt.subplots(1, 2, figsize=(11, 5.5))
    a = por_region.head(top)
    b = por_region.sort_values('patrones', ascending=False).head(top)
    axes[0].barh(range(len(a)), a['pixeles'], color=TEAL, height=0.7)
    axes[0].set_yticks(range(len(a)), a['region_id'], fontsize=9)
    axes[0].invert_yaxis()
    axes[0].xaxis.set_major_formatter(mticker.FuncFormatter(_fmt_k))
    axes[0].set_title('Más píxeles', fontsize=13, fontweight='bold', color=NAVY)
    axes[1].barh(range(len(b)), b['patrones'], color='#5B4B8A', height=0.7)
    axes[1].set_yticks(range(len(b)), b['region_id'], fontsize=9)
    axes[1].invert_yaxis()
    axes[1].set_title('Más cambios', fontsize=13, fontweight='bold', color=NAVY)
    plt.tight_layout()
    return _save(fig, '05_regiones_rank')


def fig_comunes_raras(nac, top=12):
    fig, axes = plt.subplots(1, 2, figsize=(11, 5.8))
    c = nac.head(top)
    r = nac.tail(top).sort_values('pixeles')
    axes[0].barh(range(len(c)), c['pixeles'],
                 color=[TIPO_COLOR.get(t, GRAY) for t in c['tipo']], height=0.7)
    axes[0].set_yticks(range(len(c)), c['trayectoria'], family='monospace', fontsize=8)
    axes[0].invert_yaxis()
    axes[0].xaxis.set_major_formatter(mticker.FuncFormatter(_fmt_k))
    axes[0].set_title('Comunes', fontsize=13, fontweight='bold', color=NAVY)
    axes[1].barh(range(len(r)), r['pixeles'],
                 color=[TIPO_COLOR.get(t, GRAY) for t in r['tipo']], height=0.7)
    axes[1].set_yticks(range(len(r)), r['trayectoria'], family='monospace', fontsize=8)
    axes[1].invert_yaxis()
    axes[1].set_title('Raras', fontsize=13, fontweight='bold', color=NAVY)
    plt.tight_layout()
    return _save(fig, '06_comunes_raras')


def fig_mas_menos(por_region):
    fig, axes = plt.subplots(1, 2, figsize=(11, 5))
    mas = por_region.head(8)
    menos = por_region.tail(8).sort_values('pixeles')
    axes[0].barh(range(len(mas)), mas['pixeles'], color=TEAL, height=0.7)
    axes[0].set_yticks(range(len(mas)), mas['region_id'], fontsize=9)
    axes[0].invert_yaxis()
    axes[0].xaxis.set_major_formatter(mticker.FuncFormatter(_fmt_k))
    axes[0].set_title('Más píxeles', fontsize=13, fontweight='bold', color=NAVY)
    axes[1].barh(range(len(menos)), menos['pixeles'], color=ORANGE, height=0.7)
    axes[1].set_yticks(range(len(menos)), menos['region_id'], fontsize=9)
    axes[1].invert_yaxis()
    axes[1].xaxis.set_major_formatter(mticker.FuncFormatter(_fmt_k))
    axes[1].set_title('Menos píxeles', fontsize=13, fontweight='bold', color=NAVY)
    plt.tight_layout()
    return _save(fig, '07_mas_menos')


# —— PPT helpers ——

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def add_title_slide(prs, title, subtitle):
    slide = _blank(prs)
    # fondo
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(NAVY)
    shape.line.fill.background()
    # barra accent
    bar = slide.shapes.add_shape(1, Inches(0), Inches(4.55), prs.slide_width, Inches(0.12))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(TEAL)
    bar.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.5), Inches(1.2))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = _rgb('#FFFFFF')
    p.font.name = 'Calibri'

    box2 = slide.shapes.add_textbox(Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.6))
    p2 = box2.text_frame.paragraphs[0]
    p2.text = subtitle
    p2.font.size = Pt(16)
    p2.font.color.rgb = _rgb('#A8C5C0')
    p2.font.name = 'Calibri'
    return slide


def add_section_slide(prs, title):
    slide = _blank(prs)
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(TEAL)
    shape.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.8), Inches(3.0), Inches(11.5), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(34)
    p.font.bold = True
    p.font.color.rgb = _rgb('#FFFFFF')
    p.font.name = 'Calibri'
    return slide


def add_kpi_slide(prs, kpis):
    """kpis: list of (label, value)"""
    slide = _blank(prs)
    _header(slide, 'Resumen')
    n = len(kpis)
    width = 2.6
    gap = 0.35
    total_w = n * width + (n - 1) * gap
    start = (13.333 - total_w) / 2
    for i, (lab, val) in enumerate(kpis):
        left = start + i * (width + gap)
        card = slide.shapes.add_shape(1, Inches(left), Inches(2.4), Inches(width), Inches(2.2))
        card.fill.solid()
        card.fill.fore_color.rgb = _rgb(LIGHT)
        card.line.color.rgb = _rgb('#E2E8F0')
        tb = slide.shapes.add_textbox(Inches(left + 0.15), Inches(2.7), Inches(width - 0.3), Inches(1.6))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = val
        p.font.size = Pt(26)
        p.font.bold = True
        p.font.color.rgb = _rgb(NAVY)
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = lab
        p2.font.size = Pt(12)
        p2.font.color.rgb = _rgb(GRAY)
        p2.alignment = PP_ALIGN.CENTER
    return slide


def _header(slide, title):
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(0.12), Inches(7.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(TEAL)
    bar.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.25), Inches(12), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = _rgb(NAVY)
    p.font.name = 'Calibri'


def add_chart_slide(prs, title, img_path, full=True):
    slide = _blank(prs)
    _header(slide, title)
    if full:
        slide.shapes.add_picture(str(img_path), Inches(0.5), Inches(1.0), width=Inches(12.3))
    else:
        slide.shapes.add_picture(str(img_path), Inches(1.2), Inches(1.1), width=Inches(10.8))
    return slide


def add_closing(prs):
    slide = _blank(prs)
    shape = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(NAVY)
    shape.line.fill.background()
    box = slide.shapes.add_textbox(Inches(0.8), Inches(3.1), Inches(11.5), Inches(1))
    p = box.text_frame.paragraphs[0]
    p.text = 'MapBiomas Colombia · Colección 4'
    p.font.size = Pt(22)
    p.font.color.rgb = _rgb('#FFFFFF')
    p.alignment = PP_ALIGN.CENTER
    return slide


def main():
    if not CSV.exists():
        raise FileNotFoundError(CSV)

    _style()
    df, nac, por_region, por_tipo = cargar()
    total_px = int(nac['pixeles'].sum())
    n_reg = df['region_id'].nunique()
    n_pat = len(nac)

    print('Generando figuras…')
    p_tipo = fig_tipo_donut(por_tipo, total_px)
    p_top = fig_top_tray(nac, total_px)
    p_conc, n80 = fig_concentracion(nac, total_px)
    p_ext = fig_extension(nac)
    p_reg = fig_regiones_px_cambios(por_region)
    p_cr = fig_comunes_raras(nac)
    p_mm = fig_mas_menos(por_region)

    print('Armando PPT…')
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(
        prs,
        'Trayectorias de falsos bosques',
        'Análisis nacional y por regiones · MapBiomas Colombia',
    )
    add_kpi_slide(prs, [
        ('Regiones', f'{n_reg}'),
        ('Patrones', f'{n_pat:,}'),
        ('Píxeles', f'{total_px / 1e6:.1f}M'),
        ('≈ 80% píxeles', f'{n80} patrones'),
    ])

    add_section_slide(prs, '01  ·  Nacional')
    add_chart_slide(prs, 'Tipo de anomalía', p_tipo, full=False)
    add_chart_slide(prs, 'Top trayectorias', p_top)
    add_chart_slide(prs, 'Concentración de píxeles', p_conc, full=False)
    add_chart_slide(prs, 'Extensión vs magnitud', p_ext, full=False)

    add_section_slide(prs, '02  ·  Por región')
    add_chart_slide(prs, 'Regiones: píxeles y cambios', p_reg)
    add_chart_slide(prs, 'Trayectorias comunes y raras', p_cr)
    add_chart_slide(prs, 'Regiones con más y menos píxeles', p_mm)

    add_closing(prs)

    prs.save(OUT_PPT)
    print(f'Guardado: {OUT_PPT.resolve()}')


if __name__ == '__main__':
    main()
